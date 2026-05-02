"""
Apply approved section edits back to the original .docx or .pptx.

Only text content is modified. Paragraph styles, run-level formatting
(font, size, bold, colour), images, tables, and page layout are untouched.
PDF is not patchable — the caller is responsible for skipping patch_pptx/patch_document
and instead providing only the review document for manual edits.
"""

import io
from datetime import date
from pathlib import Path

from docx import Document
from pptx import Presentation

from src.models import ChangeType, ReviewState, Section


def patch_document(
    original_bytes: bytes,
    sections: list[Section],
    original_filename: str,
) -> tuple[bytes, str]:
    """
    Return (patched_bytes, output_filename).

    For each section whose effective_text() differs from original_text,
    the corresponding paragraphs in the document are updated in-place.
    """
    doc = Document(io.BytesIO(original_bytes))
    paragraphs = doc.paragraphs

    for section in sections:
        new_text = section.effective_text()
        if new_text == section.original_text:
            continue

        new_lines = [ln for ln in new_text.split("\n") if ln.strip()]
        indices = section.paragraph_indices

        for pos, para_idx in enumerate(indices):
            if para_idx >= len(paragraphs):
                continue
            para = paragraphs[para_idx]
            line_text = new_lines[pos] if pos < len(new_lines) else ""
            _update_paragraph_text(para, line_text)

    output = io.BytesIO()
    doc.save(output)
    output.seek(0)

    stem = Path(original_filename).stem
    today = date.today().strftime("%Y-%m-%d")
    output_filename = f"{stem}_updated_{today}.docx"

    return output.getvalue(), output_filename


def _update_paragraph_text(para, new_text: str) -> None:
    """
    Write new_text into the paragraph while keeping the first run's
    character formatting (font, size, bold, italic, colour, etc.).
    Subsequent runs are cleared so no ghost text remains.
    """
    if not para.runs:
        para.add_run(new_text)
        return

    para.runs[0].text = new_text
    for run in para.runs[1:]:
        run.text = ""


def patch_pptx(
    original_bytes: bytes,
    sections: list[Section],
    original_filename: str,
) -> tuple[bytes, str]:
    """
    Return (patched_bytes, output_filename) for a .pptx file.

    Each section's paragraph_indices[0] is the slide index. Non-title body
    paragraphs on that slide are updated line-by-line from effective_text(),
    preserving the first run's character formatting.
    """
    prs = Presentation(io.BytesIO(original_bytes))

    slide_map: dict[int, Section] = {}
    for s in sections:
        if s.paragraph_indices:
            slide_map[s.paragraph_indices[0]] = s

    for slide_idx, slide in enumerate(prs.slides):
        section = slide_map.get(slide_idx)
        if not section:
            continue
        new_text = section.effective_text()
        if new_text == section.original_text:
            continue

        new_lines = [ln for ln in new_text.split("\n") if ln.strip()]

        body_paras = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if _is_title_shape(shape):
                continue
            body_paras.extend(shape.text_frame.paragraphs)

        for pos, para in enumerate(body_paras):
            line_text = new_lines[pos] if pos < len(new_lines) else ""
            _update_pptx_paragraph_text(para, line_text)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    stem = Path(original_filename).stem
    today = date.today().strftime("%Y-%m-%d")
    return output.getvalue(), f"{stem}_updated_{today}.pptx"


def _update_pptx_paragraph_text(para, new_text: str) -> None:
    """
    pptx paragraphs: add_run() takes no arguments — create the run first,
    then set its text. Subsequent runs are cleared to avoid ghost content.
    """
    if not para.runs:
        run = para.add_run()
        run.text = new_text
        return
    para.runs[0].text = new_text
    for run in para.runs[1:]:
        run.text = ""


def _is_title_shape(shape) -> bool:
    try:
        return shape.placeholder_format is not None and shape.placeholder_format.idx == 0
    except Exception:
        return False
