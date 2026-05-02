from pathlib import Path
from typing import Optional
import pdfplumber
from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation

from src.models import Section

PARAGRAPH_BLOCK_SIZE = 4
SCANNED_PDF_CHAR_THRESHOLD = 50
HEADING_STYLES = {"Heading 1", "Heading 2", "Heading 3"}


def extract_docx_sections(path: Path) -> list[Section]:
    doc = Document(str(path))
    paragraphs = doc.paragraphs

    has_headings = any(p.style and p.style.name in HEADING_STYLES for p in paragraphs if p.text.strip())
    if has_headings:
        return _split_by_headings(paragraphs)
    return _split_by_blocks(paragraphs)


def _split_by_headings(paragraphs) -> list[Section]:
    sections: list[Section] = []
    current_heading = "Introduction"
    current_text: list[str] = []
    current_indices: list[int] = []
    section_index = 0

    for i, para in enumerate(paragraphs):
        text = para.text.strip()
        if not text:
            continue

        if para.style and para.style.name in HEADING_STYLES:
            if current_text:
                sections.append(Section(
                    index=section_index,
                    heading=current_heading,
                    original_text="\n".join(current_text),
                    paragraph_indices=current_indices,
                ))
                section_index += 1
            current_heading = text
            current_text = []
            current_indices = []
        else:
            current_text.append(text)
            current_indices.append(i)

    if current_text:
        sections.append(Section(
            index=section_index,
            heading=current_heading,
            original_text="\n".join(current_text),
            paragraph_indices=current_indices,
        ))

    return sections


def _split_by_blocks(paragraphs) -> list[Section]:
    sections: list[Section] = []
    non_empty = [(i, p.text.strip()) for i, p in enumerate(paragraphs) if p.text.strip()]

    for block_start in range(0, len(non_empty), PARAGRAPH_BLOCK_SIZE):
        block = non_empty[block_start:block_start + PARAGRAPH_BLOCK_SIZE]
        indices = [i for i, _ in block]
        texts = [t for _, t in block]
        section_num = block_start // PARAGRAPH_BLOCK_SIZE + 1
        sections.append(Section(
            index=section_num - 1,
            heading=f"Section {section_num}",
            original_text="\n".join(texts),
            paragraph_indices=indices,
        ))

    return sections


def extract_pdf_text(path: Path) -> tuple[str, bool]:
    """
    Returns (extracted_text, is_scanned).
    is_scanned is True when the PDF appears to have no text layer.
    """
    text_parts: list[str] = []

    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)

    full_text = "\n\n".join(text_parts).strip()
    is_scanned = len(full_text) < SCANNED_PDF_CHAR_THRESHOLD
    return full_text, is_scanned


def extract_pptx_sections(path: Path) -> list[Section]:
    """One section per slide. Heading = slide title; body = all non-title text frames."""
    prs = Presentation(str(path))
    sections: list[Section] = []

    for slide_idx, slide in enumerate(prs.slides):
        title = f"Slide {slide_idx + 1}"
        body_lines: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if _is_title_shape(shape):
                title = text
            else:
                body_lines.append(text)

        if not body_lines:
            continue

        sections.append(Section(
            index=slide_idx,
            heading=title,
            original_text="\n".join(body_lines),
            paragraph_indices=[slide_idx],
        ))

    return sections


def extract_pdf_sections(path: Path) -> tuple[list[Section], Optional[str]]:
    """
    Extract sections from a PDF primary document using paragraph-block splitting.
    Returns (sections, error_message). error_message is None on success.
    Sections have empty paragraph_indices — PDF output is not supported.
    """
    text, is_scanned = extract_pdf_text(path)
    if is_scanned:
        return [], (
            "This PDF appears to be a scanned image with no readable text. "
            "Please provide a text-based PDF."
        )

    lines = [ln.strip() for ln in text.split("\n") if ln.strip()]
    if not lines:
        return [], "No readable text found in this PDF."

    sections: list[Section] = []
    for block_start in range(0, len(lines), PARAGRAPH_BLOCK_SIZE):
        block = lines[block_start : block_start + PARAGRAPH_BLOCK_SIZE]
        section_num = block_start // PARAGRAPH_BLOCK_SIZE + 1
        sections.append(Section(
            index=section_num - 1,
            heading=f"Section {section_num}",
            original_text="\n".join(block),
            paragraph_indices=[],
        ))

    return sections, None


def _is_title_shape(shape) -> bool:
    try:
        return shape.placeholder_format is not None and shape.placeholder_format.idx == 0
    except Exception:
        return False


def extract_reference_text(path: Path) -> tuple[Optional[str], Optional[str]]:
    """
    Returns (text, error_message). error_message is None on success.
    Supports .docx, .pdf, and .pptx reference files.
    """
    suffix = path.suffix.lower()

    if suffix == ".docx":
        doc = Document(str(path))
        text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
        return text, None

    if suffix == ".pdf":
        text, is_scanned = extract_pdf_text(path)
        if is_scanned:
            return None, (
                "This PDF appears to be a scanned image with no readable text. "
                "Please provide a text-based PDF or a Word document instead."
            )
        return text, None

    if suffix == ".pptx":
        prs = Presentation(str(path))
        lines: list[str] = []
        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        lines.append(text)
        if not lines:
            return None, "No readable text found in this PowerPoint file."
        return "\n".join(lines), None

    return None, f"Unsupported reference file type: '{suffix}'. Please upload a .docx, .pdf, or .pptx file."
